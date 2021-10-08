from __future__ import print_function
import time
import cloudmersive_convert_api_client
from cloudmersive_convert_api_client.rest import ApiException
from pprint import pprint
from pptx import Presentation
from pptx.shapes.autoshape import Shape
import copy

def summarizer(text: str):
    output = {}
    curMain = None
    for line in text.split('\n'):
        if line.startswith('  - '):
            curMain = line.replace('  - ', '')
            output[curMain] = []
        elif line.startswith('    - '):
            output[curMain].append('     ' + line)
    return output


def findTextBox(presentation: Presentation):
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if type(shape) == Shape:
                return shape.text_frame


def replace_paragraph_text_retaining_initial_formatting(paragraph, new_text):
    p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
    # remove all but the first run
    print(f'For paragraph {paragraph} with runs "{paragraph.runs}" adding text {new_text}')
    for idx, run in enumerate(paragraph.runs):
        if idx == 0:
            continue
        p.remove(run._r)
    paragraph.runs[0].text = new_text

def replaceTitle(presentation: Presentation, new_title):
    title_box = presentation.slides[0].shapes.title.text_frame.paragraphs[0]
    replace_paragraph_text_retaining_initial_formatting(title_box, new_title)




def replaceBulletPoints(textframe, inputdict: dict):

    def copy_run(run_from, run_to):
        r_to = run_to._r
        r_to.addnext(copy.deepcopy(run_from._r))
        r_to.getparent().remove(r_to)

    def copy_paragraph(paragraph_from, paragraph_to):
        p_to = paragraph_to._p
        p_to.addnext(copy.deepcopy(paragraph_from._p))
        p_to.getparent().remove(p_to)

    first_paragraph = textframe.paragraphs[0]
    first_paragraph_run = textframe.paragraphs[0].runs[0]
    i = 0
    for key in inputdict:
        if len(textframe.paragraphs) < i + 1:
            p = textframe.add_paragraph()
            copy_paragraph(paragraph_from=first_paragraph, paragraph_to=p)
        replace_paragraph_text_retaining_initial_formatting(textframe.paragraphs[i], key)
        i += 1
        for subpoints in inputdict[key]:
            p = textframe.add_paragraph()
            r = p.add_run()
            copy_run(run_from=first_paragraph_run, run_to=r)
            replace_paragraph_text_retaining_initial_formatting(textframe.paragraphs[i], subpoints)
            i += 1


def createPres(topic: list, template_slide_path, outputFolder=''):
  if 'summary' in topic:
      prs = Presentation(template_slide_path)
      replaceTitle(prs, topic['title'])
      text_box = findTextBox(prs)
      replaceBulletPoints(text_box, summarizer(topic['summary']))
      print(f"Created '{outputFolder}{topic['reference'][4:-1]}.pptx'")
      prs.save(f"{outputFolder}{topic['reference'][4:-1]}.pptx")
      return prs
  else:
      print(f"Topic {topic['title']} has no sumamry. Skipping PPT creation.")
      return None

#topic = {
#    "reference": ".. _data-driven-architecture-1-1:",
#    "title": "Data Driven Architecture (DDA)",
#    "source_file": "training/010_architecture_and_concepts/020_Beginner.rst",
#    "text": ".. _data-driven-architecture-1-1:\n\nData Driven Architecture (DDA)\n^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^\n\n.. sidebar:: Summary\n\n  - Asynchronous data transfer\n  - Data centric\n  - Very scaleable\n  - Agile\n\nAs opposed to both the P2P and the ESB integration principles the Data\nDriven Architecture (DDA) does not focus on systems but rather the data\nthese systems store and how it can be used in a data-centric ecosystem.\nThis gives us an agile, robust IA. See figure :ref:`figure-dda-1-1`.\n\n.. _figure-dda-1-1:\n.. figure:: ./media/Data_Driven_Architecture.png\n   :align: center\n\n   Data Driven Architecture\n\nAs opposed to P2P and ESB, DDA is, respectively, scalable and agile - in\nthat it does not need the \u201cBus\u201d to orchestrate data flows, rather DDA\nrelies on retrieving all the data in a system and connecting it internally\nto enhance and propagate it for usage in outbound flows.\n\nAs in all great things, there is risk involved.\nIn order to utilize DDA effectively you need logical and robust principles\nto create flexible data flows and models.\nThis can be achieved by always thinking ahead and leaving room for growth,\nwhich you will learn how to do throughout this course.\nIf done correctly nothing beats the scalability, cost-effectiveness and\npossibilities a DDA provides.\n\n.. seealso::\n\n  TODO\n\n",
#    "summary": "  - Asynchronous data transfer\n  - Data centric\n  - Very scaleable\n  - Agile\n"
#}
topics = [{
  "reference": ".. _environment-variables-secrets-2-1:",
  "title": "Environment variables & Secrets",
  "source_file": "training/020_systems/020_Beginner.rst",
  "text": ".. _environment-variables-secrets-2-1:\n\nEnvironment variables & Secrets\n~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~\n\n.. sidebar:: Summary\n\n  - Environment variables and secrets are named values used to parameterize configs\n  - Environment variables are:\n\n    - unencrypted\n    - referenced with: ``\"$ENV(my-env-var)\"``\n\n  - Secrets are:\n\n    - encrypted\n    - referenced with: ``\"$SECRET(my-secret)\"``\n\n  - Both are defined under **Datahub > Variables**\n  - Secrets can also be defined under a system's **Secrets** tab\n  - Eases and improves config maintenance\n\nIn this section we will cover how environment variables and secrets typically\nare used in system configs.\n\nEnvironment variables and secrets are named values\nthat can be used to parameterize Sesam configs.\n\nEnvironment variables are stored and processed as *unencrypted* values,\nand are referenced with ``\"$ENV(my-env-var)\"``.\n\nSecrets are stored and processed as *encrypted* values,\nand are referenced with ``\"$SECRET(my-secret)\"``.\n\nBoth are defined in the Sesam Management Studio under **Datahub > Variables**.\n\nSecrets can also be defined locally in a system config under the system's\n**Secrets** tab.\n\n.. warning::\n\n  If a system config is deleted, all secrets stored locally in that system config is lost!\n\nIt is generally a good idea to put the parts of a configuration that differ between\nenvironments (develop, test, production, etc.) into environment variables.\nThis includes configs such as server names, database connection strings, API URLs, usernames, etc.\n\nBy putting these config parts into environment variables you can define each of them\nseparately in each Sesam node used for the respective environments,\nbut keep the actual system config identical in each node.\n\nThis is also practical for version control of the config.\nYou can change the values of the environment variables separate from the actual\nsystem config.\n\nContinuing from the example :ref:`practice-system-config-initial`, let us see how the\nintroduction of environment variables can improve the system config.\nThe ``url_pattern`` is a good canditate to be put into an environment variable.\nLet us call it `\"difi-api\"` and reference it from the system config.\n\nFirst we define the new environment variable under\n**Datahub > Variables > Environment variables**:\n\n.. code-block:: json\n\n  \"difi-api\": \"https://ws.geonorge.no/kommuneinfo/v1/%s\"\n\nThen we change the system config to reference it:\n\n.. _practice-system-config-env-var-ref:\n.. code-block:: json\n  :caption: Practice system config with environment variable reference\n  :linenos:\n\n  {\n    \"_id\": \"difi\",\n    \"type\": \"system:url\",\n    \"url_pattern\": \"$ENV(difi-api)\",\n    \"verify_ssl\": true\n  }\n\nSay we want to access different Difi APIs depending on which environment\nwe are accessing Difi from, or that Difi decided to change the API URL at some point.\nThe only thing that we have to update is the value of the ``difi-api``\nenvironment variable.\nNo changes to the actual system config is required.\n\n.. seealso::\n\n  Concepts - Configuration: :ref:`concepts-environment-variables`\n\n  Concepts - Configuration: :ref:`concepts-secrets`\n\n",
  "summary": "  - Environment variables and secrets are named values used to parameterize configs\n  - Environment variables are:\n    - unencrypted\n    - referenced with: ``\"$ENV(my-env-var)\"``\n  - Secrets are:\n    - encrypted\n    - referenced with: ``\"$SECRET(my-secret)\"``\n  - Both are defined under **Datahub > Variables**\n  - Secrets can also be defined under a system's **Secrets** tab\n  - Eases and improves config maintenance\n"
},
{
  "reference": ".. _systems-1-1:",
  "title": "The MSSQL system",
  "source_file": "training/010_architecture_and_concepts/020_Beginner.rst",
  "text": ".. _systems-1-1:\n\nSystems\n~~~~~~~\n\n.. sidebar:: Summary\n\n  Systems are interfaces to external systems.\n\nSystems are one of Sesam's core components.\nSystems can connect to external providers such as SQL databases, REST APIs,\nMicroservices and more, to either import data into Sesam or export data out from Sesam.\nSystems are therefore the start and end points of every dataflow.\n\nSystems may cover other functionalities as well, but we will cover those special cases\nlater.\n\nIn this section we will show you an example of the most commom system in a Sesam installation,\nthe mssql system. We will also show how this system can connect to pipes to\neither import or export data, depending on your need.\n\nThe MSSQL system\n^^^^^^^^^^^^^^^^\n\n.. figure:: ./media/mssql-system-config.png\n   :align: right\n   :alt: MSSQL system config.\n\n   MSSQL system config\n\nSince they are a relatively common way to store data, Sesam has a ready built-in connector for MSSQL databases. The MSSQL system inside Sesam connects to an MSSQL database by sending the host, database and port information, as well as authentication parameters, through a built in connector inside Sesam. Note that in the system config we also have to specify the system type ``system:mssql``.\n\n.. figure:: ./media/mssql-system-status.png\n   :align: right\n   :alt: MSSQL system status.\n\n   MSSQL system status.\n\nOnce the connection is open the node can extract data from the tables in the database through inbound pipes connected to the system. You can see if the connection to the MSSQL database is open by going to the \"Status\" tab on the system page. Should the system health state \"failure\" in your connectivity, this could be because you have some parameter values in your config wrong, or there might be a firewall blocking your access.\n\n.. seealso::\n\n  Learn Sesam: :ref:`systems-beginner-2-1`\n\n  Developer Guide - Service Configuration: :ref:`system_section`\n\n",
  "summary": "  Systems are interfaces to external systems.\n"
}]
#template_slide_path = 'splitme.pptx'
#for t in topics:
#    curPres = createPres(topic, template_slide_path)



# Configure API key authorization: Apikey
configuration = cloudmersive_convert_api_client.Configuration()
configuration.api_key['Apikey'] = 'dac7e8ea-4a89-49e3-8da7-646fb9186784'
# Uncomment below to setup prefix (e.g. Bearer) for API key, if needed
#configuration.api_key_prefix['Apikey'] = 'Bearer'
# create an instance of the API class
api_instance = cloudmersive_convert_api_client.MergeDocumentApi(cloudmersive_convert_api_client.ApiClient(configuration))
input_file1 = '/Users/gabriell.vig/Work/Sesam/docs/training/courses/tc_data-driven-architecture-1-1.pptx' # file | First input file to perform the operation on.
input_file2 = '/Users/gabriell.vig/Work/Sesam/docs/training/courses/tc_datasets-1-1.pptx' # file | Second input file to perform the operation on (more than 2 can be supplied).
try:
    # Merge Two PowerPoint PPTX Together
    api_response = api_instance.merge_document_pptx(input_file1, input_file2)
    pprint(api_response)
    with open("output.pptx", "wb") as binary_file:
        binary_file.write(api_response)
except ApiException as e:
    print("Exception when calling MergeDocumentApi->merge_document_pptx: %s\n" % e)
