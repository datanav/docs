from pptx import Presentation
from pptx.shapes.autoshape import Shape
import copy

prs = Presentation('splitme.pptx')
#print(prs)
text_runs = []

text_box = None
for slide in prs.slides:
    for shape in slide.shapes:
        print('#')
        print(type(shape))
        print(shape)
        if not shape.has_text_frame:
            continue
        if type(shape) == Shape:
            text_box = shape
        for paragraph in shape.text_frame.paragraphs:
            print('##')
            print(paragraph)
            for run in paragraph.runs:
                print('###')
                print(run)
                print(run.text)
                pass

from sys import exit
#exit(0)

def replace_paragraph_text_retaining_initial_formatting(paragraph, new_text):
    p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
    # remove all but the first run
    print(f'For paragraph {paragraph} with runs "{paragraph.runs}" adding text {new_text}')
    for idx, run in enumerate(paragraph.runs):
        if idx == 0:
            continue
        p.remove(run._r)
    paragraph.runs[0].text = new_text

def copy_run(run_from, run_to):
    r_to = run_to._r
    r_to.addnext(copy.deepcopy(run_from._r))
    r_to.getparent().remove(r_to)

def copy_paragraph(paragraph_from, paragraph_to):
    p_to = paragraph_to._p
    p_to.addnext(copy.deepcopy(paragraph_from._p))
    p_to.getparent().remove(p_to)

def replaceBulletPoints(textframe, text):
    first_paragraph = textframe.paragraphs[0]
    #for i, line in enumerate(text.split('\n')):
    for i, line in enumerate(text):
        if len(textframe.paragraphs) < i + 1:
            p = textframe.add_paragraph()
            copy_paragraph(paragraph_from=first_paragraph, paragraph_to=p)
#            p.level = 3
            #print(p.indent)
        print(len(textframe.paragraphs))
        print(i + 1)
        replace_paragraph_text_retaining_initial_formatting(textframe.paragraphs[i], line)

def replaceBulletPoints2(textframe, inputdict: dict):
    first_paragraph = textframe.paragraphs[0]
    first_paragraph_run = textframe.paragraphs[0].runs[0]
    i = 0
    for key in inputdict:
        if len(textframe.paragraphs) < i + 1:
            p = textframe.add_paragraph()
            copy_paragraph(paragraph_from=first_paragraph, paragraph_to=p)
        replace_paragraph_text_retaining_initial_formatting(textframe.paragraphs[i], key)
        i += 1
        for subpoints in inputdict[key]:
            p = textframe.add_paragraph()
            r = p.add_run()
            copy_run(run_from=first_paragraph_run, run_to=r)
            replace_paragraph_text_retaining_initial_formatting(textframe.paragraphs[i], subpoints)
            i += 1




writethis = """  - Environment variables and secrets are named values used to parameterize configs
  - Environment variables are:
    - unencrypted
    - referenced with: ``"$ENV(my-env-var)"``
  - Secrets are:
    - encrypted
    - referenced with: ``"$SECRET(my-secret)"``
  - Both are defined under **Datahub > Variables**
  - Secrets can also be defined under a system's **Secrets** tab
  - Eases and improves config maintenance
"""

def summarizer(text: str):
    output = {}
    curMain = None
    for line in text.split('\n'):
        if line.startswith('  - '):
            curMain = line.replace('  - ', '')
            output[curMain] = []
        elif line.startswith('    - '):
            output[curMain].append(line)
    return output

def summarizer2(text: str):
    output = []
    curMain = None
    indexer = 0
    for line in text.split('\n'):
        if line.startswith('  - '):
            if len(output) > 0:
                indexer += 1
            output.append(line + '\n')
        elif line.startswith('    - '):
            output[indexer] += line + '\n'
    return output

from json import dumps

s = prs.slides[0]

replace_paragraph_text_retaining_initial_formatting(s.shapes.title.text_frame.paragraphs[0], 'Environment variables & Secrets')
#replaceBulletPoints(text_box.text_frame, writethis)
#replaceBulletPoints(text_box.text_frame, writethis.split('\n  - '))
#replaceBulletPoints(text_box.text_frame, summarizer2(writethis))
replaceBulletPoints2(text_box.text_frame, summarizer(writethis))

prs.save('tmp_output.pptx')
